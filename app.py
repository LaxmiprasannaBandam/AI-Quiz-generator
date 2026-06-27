import os
from flask import Flask, render_template, request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from pptx import Presentation

app = Flask(__name__)

# Configuration
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev-secret-key-change-in-production')
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16 MB max upload
app.config['ALLOWED_EXTENSIONS'] = {'pptx'}

# Ensure upload directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)


def allowed_file(filename):
    """Check if the uploaded file has an allowed extension."""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']


def extract_pptx_text(filepath):
    """
    Extract all text from a PPTX file.
    Returns a list of dicts, one per slide, each containing:
      - slide_number: int
      - text: str (all text content from that slide)
    """
    slides_data = []
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)
            slides_data.append({
                'slide_number': i,
                'text': '\n'.join(slide_text_parts)
            })
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from PPTX: {str(e)}")
    return slides_data


@app.route('/')
def index():
    """Render the main page."""
    return render_template('index.html')


@app.route('/api/health', methods=['GET'])
def health_check():
    """Health check endpoint."""
    return jsonify({'status': 'healthy', 'message': 'AI Quiz Generator is running'})


@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Handle file upload for quiz generation."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        return jsonify({
            'message': 'File uploaded successfully',
            'filename': filename,
            'filepath': filepath
        }), 200

    return jsonify({'error': 'File type not allowed. Only .pptx files are accepted.'}), 400


@app.route('/api/extract-text', methods=['POST'])
def extract_text():
    """
    Extract text from an uploaded PPTX file.
    Expects JSON: { "filename": "uploaded_file.pptx" }
    Returns JSON with slides array containing slide_number and text.
    """
    data = request.get_json()
    if not data or 'filename' not in data:
        return jsonify({'error': 'No filename provided'}), 400

    filename = data['filename']
    # Security: ensure the filename is safe and exists in our uploads directory
    safe_filename = secure_filename(filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], safe_filename)

    if not os.path.exists(filepath):
        return jsonify({'error': 'File not found on server'}), 404

    try:
        slides = extract_pptx_text(filepath)
        return jsonify({
            'filename': safe_filename,
            'total_slides': len(slides),
            'slides': slides
        }), 200
    except RuntimeError as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate-quiz', methods=['POST'])
def generate_quiz():
    """Generate quiz questions from uploaded content."""
    data = request.get_json()

    if not data or 'content' not in data:
        return jsonify({'error': 'No content provided'}), 400

    # Placeholder for AI quiz generation logic
    # This will be implemented with AI integration later
    content = data['content']
    num_questions = data.get('num_questions', 5)

    # Mock response for now
    quiz = {
        'title': 'Generated Quiz',
        'questions': [
            {
                'id': 1,
                'question': f'Sample question based on: {content[:50]}...',
                'options': ['Option A', 'Option B', 'Option C', 'Option D'],
                'correct_answer': 'Option A'
            }
        ],
        'total_questions': num_questions
    }

    return jsonify(quiz), 200


@app.errorhandler(404)
def not_found(error):
    """Handle 404 errors."""
    return jsonify({'error': 'Not found'}), 404


@app.errorhandler(500)
def internal_error(error):
    """Handle 500 errors."""
    return jsonify({'error': 'Internal server error'}), 500


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)